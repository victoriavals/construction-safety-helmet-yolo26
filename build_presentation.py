# -*- coding: utf-8 -*-
"""Generate PPTX presentasi (sidang/akademik) + script presentasi — versi DETAIL.
Pakai data nyata terkini (runs/ + presentation_figures/). Jalankan via Python GLOBAL.
"""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

ROOT = Path(__file__).resolve().parent
OUT = ROOT / "presentasi"; OUT.mkdir(exist_ok=True)
NAVY = RGBColor(0x1F, 0x3A, 0x5F); BLUE = RGBColor(0x4C, 0x72, 0xB0)
RED = RGBColor(0xC4, 0x4E, 0x52); GREEN = RGBColor(0x55, 0xA8, 0x68)
GREY = RGBColor(0x55, 0x55, 0x55); LGREY = RGBColor(0xEE, 0xEE, 0xEE); WHITE = RGBColor(0xFF, 0xFF, 0xFF)

prs = Presentation(); prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]; SW, SH = prs.slide_width, prs.slide_height
SLIDES = []


def _tf(slide, l, t, w, h):
    tb = slide.shapes.add_textbox(l, t, w, h); tb.text_frame.word_wrap = True; return tb.text_frame


def title_bar(slide, title, num):
    bar = slide.shapes.add_shape(1, 0, 0, SW, Inches(1.0))
    bar.fill.solid(); bar.fill.fore_color.rgb = NAVY; bar.line.fill.background()
    tf = bar.text_frame; tf.word_wrap = True; tf.margin_left = Inches(0.35); tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.text = title; p.font.size = Pt(24); p.font.bold = True; p.font.color.rgb = WHITE
    nt = _tf(slide, SW - Inches(1.0), Inches(0.3), Inches(0.8), Inches(0.4))
    pp = nt.paragraphs[0]; pp.text = str(num); pp.alignment = PP_ALIGN.RIGHT; pp.font.size = Pt(13); pp.font.color.rgb = WHITE


def bullets(tf, items):
    for i, b in enumerate(items):
        lvl = 0
        if isinstance(b, tuple): b, lvl = b
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = ("• " if lvl == 0 else "– ") + b; p.level = lvl
        p.font.size = Pt(16 if lvl == 0 else 14); p.font.color.rgb = NAVY if lvl == 0 else GREY
        p.space_after = Pt(6)


def add_image(slide, relpath, l, t, max_w, max_h):
    p = ROOT / relpath
    if not p.exists():
        box = slide.shapes.add_shape(1, l, t, max_w, max_h)
        box.fill.solid(); box.fill.fore_color.rgb = LGREY; box.line.color.rgb = RED
        tfx = box.text_frame; tfx.word_wrap = True; tfx.vertical_anchor = MSO_ANCHOR.MIDDLE
        pr = tfx.paragraphs[0]; pr.text = "GANTI: " + relpath; pr.alignment = PP_ALIGN.CENTER
        pr.font.size = Pt(13); pr.font.color.rgb = RED; pr.font.bold = True
        return
    pic = slide.shapes.add_picture(str(p), l, t, width=max_w)
    if pic.height > max_h:
        f = max_h / pic.height; pic.width = int(pic.width * f); pic.height = int(pic.height * f)
    pic.left = int(l + (max_w - pic.width) / 2)


def notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


def text_slide(num, title, items, note):
    s = prs.slides.add_slide(BLANK); title_bar(s, title, num)
    bullets(_tf(s, Inches(0.7), Inches(1.3), Inches(12), Inches(5.8)), items); notes(s, note)
    SLIDES.append((num, title, items, note))


def figure_slide(num, title, items, img, caption, note):
    s = prs.slides.add_slide(BLANK); title_bar(s, title, num)
    bullets(_tf(s, Inches(0.45), Inches(1.25), Inches(4.9), Inches(5.8)), items)
    add_image(s, img, Inches(5.6), Inches(1.25), Inches(7.3), Inches(5.1))
    cf = _tf(s, Inches(5.6), Inches(6.5), Inches(7.3), Inches(0.5))
    cp = cf.paragraphs[0]; cp.text = caption; cp.alignment = PP_ALIGN.CENTER
    cp.font.size = Pt(12); cp.font.italic = True; cp.font.color.rgb = GREY
    notes(s, note); SLIDES.append((num, title, items + [f"[GAMBAR: {caption}]"], note))


def table_slide(num, title, headers, rows, note, widths=None, intro=None):
    s = prs.slides.add_slide(BLANK); title_bar(s, title, num)
    top = Inches(1.2)
    if intro:
        itf = _tf(s, Inches(0.5), Inches(1.1), Inches(12.3), Inches(0.5))
        ip = itf.paragraphs[0]; ip.text = intro; ip.font.size = Pt(13); ip.font.italic = True; ip.font.color.rgb = GREY
        top = Inches(1.65)
    nr, nc = len(rows) + 1, len(headers)
    gt = s.shapes.add_table(nr, nc, Inches(0.5), top, Inches(12.3), Inches(0.4 * nr)).table
    if widths:
        for j, w in enumerate(widths): gt.columns[j].width = Inches(w)
    for j, h in enumerate(headers):
        c = gt.cell(0, j); c.text = h
        c.fill.solid(); c.fill.fore_color.rgb = NAVY
        pr = c.text_frame.paragraphs[0]; pr.font.size = Pt(13); pr.font.bold = True; pr.font.color.rgb = WHITE
    for i, row in enumerate(rows, 1):
        for j, val in enumerate(row):
            c = gt.cell(i, j); c.text = str(val)
            c.fill.solid(); c.fill.fore_color.rgb = WHITE if i % 2 else LGREY
            for pp in c.text_frame.paragraphs:
                pp.font.size = Pt(11.5); pp.font.color.rgb = NAVY if j == 0 else GREY
                if j == 0: pp.font.bold = True
    notes(s, note)
    SLIDES.append((num, title, [f"{r[0]}: {r[1]}" for r in rows], note))


def cover(title, subtitle, lines, note):
    s = prs.slides.add_slide(BLANK)
    bg = s.shapes.add_shape(1, 0, 0, SW, SH); bg.fill.solid(); bg.fill.fore_color.rgb = NAVY; bg.line.fill.background()
    tf = _tf(s, Inches(1), Inches(2.2), Inches(11.3), Inches(2.4))
    p = tf.paragraphs[0]; p.text = title; p.font.size = Pt(40); p.font.bold = True; p.font.color.rgb = WHITE
    p2 = tf.add_paragraph(); p2.text = subtitle; p2.font.size = Pt(22); p2.font.color.rgb = RGBColor(0xBB, 0xCC, 0xEE)
    tf3 = _tf(s, Inches(1), Inches(5.0), Inches(11.3), Inches(1.8))
    for i, line in enumerate(lines):
        pp = tf3.paragraphs[0] if i == 0 else tf3.add_paragraph()
        pp.text = line; pp.font.size = Pt(16); pp.font.color.rgb = WHITE
    notes(s, note); SLIDES.append((len(SLIDES) + 1, title, lines, note))


# ===================== METRIK NYATA (dari runs/) ======================= #
# Training (best @ epoch 40): P R mAP50 mAP50-95
TR = dict(P=0.848, R=0.708, m50=0.786, m5095=0.508)
VAL = dict(P=0.845, R=0.711, m50=0.687, m5095=0.453)   # eval conf=0.25
TST = dict(P=0.702, R=0.751, m50=0.648, m5095=0.405)

# ============================ SLIDES =================================== #
cover("Deteksi Keselamatan Helm Konstruksi", "Object Detection 3 Kelas Berbasis YOLO26s",
      ["[Nama Anda]  •  [NIM]", "[Program Studi / Instansi]", "[Tanggal Sidang]"],
      "Pembuka: perkenalkan diri & judul. Proyek: pipeline lengkap deteksi helm K3 (Helmet/No-Helmet/Person) dengan YOLO26s, dari analisis data hingga model siap deploy + demo web.")

text_slide(2, "Latar Belakang & Masalah", [
    "K3 konstruksi: helm wajib, tapi kepatuhan sulit dipantau manual.",
    "Pelanggaran 'tidak memakai helm' berisiko cedera kepala fatal.",
    "Pengawasan CCTV manual: lelah, lambat, tidak konsisten.",
    "Solusi: deteksi otomatis berbasis computer vision (object detection).",
    ("3 kelas: Helmet (pakai) · No-Helmet (tidak pakai) · Person.", 1),
], "Konteks K3: helm krusial tapi pengawasan manual tak efektif. Computer vision memantau otomatis & real-time. Kelas No-Helmet paling penting (menandai pelanggaran).")

text_slide(3, "Tujuan & Kontribusi", [
    "Membangun pipeline deteksi helm 3 kelas yang lengkap & reproducible.",
    "Tahap: EDA → pre-processing → training → tuning → evaluasi → inference → export.",
    "Efisien di GPU terbatas (RTX 4060 Ti, 8 GB VRAM).",
    "Model siap-deploy (ONNX) + demo web + tiap keputusan dijustifikasi data.",
], "Tujuan: pipeline utuh & terdokumentasi, bukan sekadar melatih model. Tantangan: GPU 8 GB. Setiap keputusan teknis berbasis data.")

text_slide(4, "Dataset", [
    "Sumber: Roboflow 'construction-safety-helmet' (CC BY 4.0).",
    "Total 19.456 gambar — train 15.555 / valid 1.945 / test 1.945.",
    "Format YOLO: satu .txt per gambar (class cx cy w h ternormalisasi).",
    "3 kelas: 0=Helmet, 1=No-Helmet, 2=Person.",
    ("Roboflow: 'No pre-processing or augmentation applied' → data MENTAH.", 1),
], "Sumber, ukuran, format YOLO, 3 kelas. Penting: data mentah (ukuran beragam, tanpa augmentasi Roboflow) — relevan untuk slide pre-processing.")

text_slide(5, "Metodologi — Pipeline", [
    "8 tahap modular (script src/ + notebook visual):",
    ("Cek GPU → EDA → Pre-processing (runtime) → Training", 1),
    ("→ Tuning → Evaluasi → Inference → Export (ONNX)", 1),
    "Konfigurasi terpusat (config.yaml); output ke runs/.",
    "Ketahanan: fallback OOM otomatis batch 8→4→2, logging, exit-code.",
], "Pipeline sebagai metodologi: tiap tahap = script + notebook. Fitur rekayasa: konfigurasi terpusat & fallback OOM agar tetap jalan di 8 GB.")

text_slide(6, "Mengapa YOLO26s?", [
    "YOLO = detektor satu-tahap: cepat, cocok real-time/CCTV.",
    "Varian 'small' (yolo26s): ~9,4 juta parameter, ~20,5 GFLOPs.",
    "Seimbang akurasi vs beban komputasi; muat di VRAM 8 GB.",
    "Kecepatan inferensi ~2,8 ms/gambar di GPU (~350+ FPS).",
], "Justifikasi model: one-stage cepat (real-time), varian small agar muat 8 GB & tetap cepat. Sebutkan jumlah parameter & kecepatan.")

# ---- EDA ----
figure_slide(7, "EDA — Distribusi Kelas", [
    "No-Helmet adalah kelas MINORITAS.",
    "Imbalance → recall kelas minoritas berisiko rendah.",
    "Padahal No-Helmet paling penting (pelanggaran K3).",
], "presentation_figures/02_eda_01.png", "Distribusi bounding box per kelas",
   "Ketimpangan kelas: No-Helmet paling sedikit, padahal terpenting untuk K3 → benang merah ke recall rendah di evaluasi.")
figure_slide(8, "EDA — Ukuran Objek (kenapa imgsz=640)", [
    "Porsi objek KECIL signifikan (small < 32²px @640).",
    "Objek kecil → resolusi rendah menghapus detail.",
    "Justifikasi: pertahankan imgsz=640.",
], "presentation_figures/02_eda_03.png", "Distribusi area bbox & komposisi small/medium/large",
   "Justifikasi data untuk imgsz=640: banyak objek kecil; menurunkan resolusi menghilangkan mereka.")
figure_slide(9, "EDA — Posisi Objek (Heatmap)", [
    "Objek terkonsentrasi di area tertentu (helm di atas).",
    "Model mempelajari bias posisi ini.",
    "Memandu augmentasi yang aman (crop/flip).",
], "presentation_figures/02_eda_06.png", "Peta panas posisi objek (overall & per kelas)",
   "Bias spasial: helm di atas, orang di tengah. Memandu pilihan augmentasi (mis. tak flip vertikal).")
figure_slide(10, "EDA — Kepadatan Objek", [
    "Sebagian gambar padat → kotak bertumpuk.",
    "→ ambang NMS (iou) berpengaruh.",
    "Gambar tanpa anotasi = contoh background.",
], "presentation_figures/02_eda_07.png", "Distribusi jumlah objek per gambar",
   "Gambar padat → ambang NMS penting. Gambar background menekan false positive.")
figure_slide(11, "EDA — Kualitas Anotasi", [
    "Label nyasar, objek sangat kecil, objek terpotong, dugaan duplikat.",
    "Noise anotasi menambah kesulitan; sebagian dibuang pipeline.",
], "presentation_figures/02_eda_09.png", "Diagnostik kualitas anotasi",
   "Kejujuran data: ada label nyasar/terpotong/duplikat. Pipeline membuang yang invalid. Menjelaskan sebagian error.")
figure_slide(12, "EDA — Kebocoran Data antar-Split", [
    "Uji near-duplicate (perceptual hash + Hamming).",
    "Kebocoran → metrik val/test terlalu optimis.",
    "Diagnostik penting untuk validitas evaluasi.",
], "presentation_figures/02_eda_16.png", "Jarak Hamming minimum val/test ke train",
   "Cek kebocoran data antar-split untuk validitas evaluasi — menunjukkan kesadaran metodologis.")

text_slide(13, "Pre-processing (Otomatis, Runtime)", [
    "Dataset mentah → TIDAK ada script pre-processing terpisah (sengaja).",
    "Ditangani OTOMATIS oleh Ultralytics saat training/inferensi:",
    ("1. Letterbox: resize jaga rasio + padding ke 640×640.", 1),
    ("2. Normalisasi piksel 0–255 → 0–1.", 1),
    ("3. Augmentasi (saat training) — lihat slide berikut.", 1),
    "Alasan: dataset read-only & Ultralytics menangani sendiri (manual = redundan/berisiko).",
], "Pre-processing ada tapi otomatis di runtime (letterbox+normalisasi+augmentasi). Manual = redundan & berisiko merusak objek kecil.")

# ================== TRAINING (DETAIL) ================== #
table_slide(14, "Training — Parameter & Alasan",
    ["Parameter", "Nilai", "Mengapa dipilih"],
    [["model", "yolo26s.pt", "Small: muat 8 GB, cepat, akurasi cukup"],
     ["imgsz", "640", "Banyak objek kecil (EDA) → resolusi memadai"],
     ["epochs", "50", "Cukup konvergen untuk dataset ini"],
     ["patience", "15", "Early-stopping: berhenti bila tak membaik (cegah overfit, hemat waktu)"],
     ["batch", "8", "Batas VRAM 8 GB @640; fallback OOM 8→4→2"],
     ["amp", "True", "Mixed precision: hemat VRAM & lebih cepat"],
     ["optimizer", "auto", "Ultralytics pilih AdamW + lr otomatis (praktis) — catatan: mengabaikan lr0"],
     ["lr0 / lrf", "0.01 / 0.01", "Lr awal & faktor lr akhir (penjadwalan menurun)"],
     ["momentum / weight_decay", "0.937 / 0.0005", "Default YOLO; weight_decay = regularisasi cegah overfit"],
     ["warmup_epochs", "3", "Lr naik bertahap di awal → stabilkan training"],
     ["seed", "42", "Reproducibility (hasil dapat diulang)"]],
    "Jelaskan tiap parameter & alasannya. Tekankan: imgsz 640 (objek kecil), batch 8 + fallback OOM (VRAM 8 GB), patience 15 (anti-overfit), AMP (efisiensi). optimizer='auto' penting — akan dibahas di tuning.",
    widths=[3.0, 1.8, 7.5])

table_slide(15, "Training — Augmentasi & Alasan",
    ["Augmentasi", "Nilai", "Mengapa"],
    [["mosaic", "1.0", "Gabung 4 gambar → variasi konteks & skala (bagus utk objek kecil)"],
     ["close_mosaic", "10", "Matikan mosaic 10 epoch terakhir → fine-tune pada gambar asli"],
     ["hsv_h/s/v", "0.015/0.7/0.4", "Jitter warna → tahan variasi pencahayaan (data blur/low-light)"],
     ["fliplr", "0.5", "Flip horizontal 50% → objek simetris kiri-kanan, aman"],
     ["flipud / degrees", "0 / 0", "TIDAK flip vertikal/rotasi → orang & helm tak terbalik di dunia nyata"],
     ["scale / translate", "0.5 / 0.1", "Variasi skala & posisi → robust ukuran/posisi objek"],
     ["erasing", "0.4", "Random erasing → robust terhadap oklusi"],
     ["mixup", "0.0", "Dimatikan (dapat membingungkan untuk deteksi)"]],
    "Augmentasi dijustifikasi konteks: mosaic & scale untuk objek kecil/variasi; HSV untuk low-light; flip horizontal aman tapi vertikal/rotasi dimatikan (tak realistis); erasing untuk oklusi.",
    widths=[3.0, 2.0, 7.3])

figure_slide(16, "Training — Preview Augmentasi", [
    "Batch nyata yang 'dilihat' model setelah letterbox + mosaic.",
    "Gabungan 4 gambar, jitter warna, flip — bukan gambar asli.",
    "Memperkaya variasi data tanpa menambah anotasi.",
], "runs/train/helmet_yolo26s_baseline/train_batch0.jpg", "train_batch0.jpg — batch training ter-augmentasi",
   "Tunjukkan wujud nyata data setelah augmentasi (mosaic). Inilah yang masuk ke model, menjelaskan slide augmentasi sebelumnya.")

figure_slide(17, "Training — Kurva Proses (per Epoch)", [
    "Loss (box/cls/dfl) train & val turun → model belajar.",
    "Metrik (P/R/mAP) naik & stabil menjelang akhir.",
    "Best model di epoch 40 (early-stopping patience=15).",
], "runs/train/helmet_yolo26s_baseline/results.png", "results.png — kurva loss & metrik per epoch",
   "Dashboard Ultralytics: loss menurun, metrik meningkat. Best di epoch 40 → early stopping bekerja. Tunjukkan tidak ada divergensi (training sehat).")

figure_slide(18, "Training — Hasil & Confusion Matrix", [
    f"Precision {TR['P']:.3f} · Recall {TR['R']:.3f}",
    f"mAP50 {TR['m50']:.3f} · mAP50-95 {TR['m5095']:.3f}",
    "Pola: Precision ≫ Recall → model KONSERVATIF.",
    ("Deteksi akurat, sebagian objek terlewat (recall).", 1),
    "Confusion: kelas mana paling tertukar/terlewat.",
], "runs/train/helmet_yolo26s_baseline/confusion_matrix_normalized.png",
   "Confusion matrix (ternormalisasi) — validasi training",
   "Metrik final (best@40). Precision tinggi, recall moderat = konservatif. Confusion matrix menunjukkan kelas mana tertukar (mis. No-Helmet ↔ Helmet/Person). Kaitkan ke imbalance EDA.")

figure_slide(19, "Training — Prediksi pada Batch Validasi", [
    "Kiri-kanan: ground-truth vs prediksi model.",
    "Kualitatif: model mengenali helm & orang dengan baik.",
    "Objek kecil/oklusi: kadang terlewat (recall).",
], "runs/train/helmet_yolo26s_baseline/val_batch0_pred.jpg", "val_batch0_pred.jpg — prediksi model pada validasi",
   "Bukti kualitatif model belajar: prediksi pada batch validasi. Bandingkan dengan val_batch0_labels.jpg (GT). Soroti objek yang terlewat = kasus sulit.")

# ================== TUNING (DETAIL) ================== #
table_slide(20, "Tuning — Rancangan & Alasan",
    ["Aspek", "Nilai", "Mengapa"],
    [["Jumlah eksperimen", "maks 4 (terkurasi)", "Ringan & aman untuk GPU 8 GB (bukan grid search berat)"],
     ["Variasi lr0", "0.01 / 0.005 / 0.003", "Eksplorasi learning rate"],
     ["Variasi epochs", "50 / 40 / 30", "Cek pengaruh durasi training"],
     ["Variasi weight_decay", "0.0005 / 0.0001", "Cek pengaruh regularisasi"],
     ["Model", "hanya small", "Lindungi VRAM 8 GB (medium/large ditolak)"],
     ["Fallback OOM", "batch 8→4→2", "Tetap jalan bila kehabisan memori"],
     ["Kriteria pemenang", "mAP50-95 → recall → stabilitas", "Utamakan akurasi lokalisasi, lalu sensitivitas, lalu kestabilan"]],
    "Tuning sengaja ringan & terkurasi (4 eksperimen) demi VRAM 8 GB. Variasi pada lr0, epochs, weight_decay. Pemenang dipilih berlapis: mAP50-95 dulu, lalu recall, lalu stabilitas.",
    widths=[3.2, 2.6, 6.5])

figure_slide(21, "Tuning — Hasil 4 Eksperimen", [
    "exp_01 (lr0.01, 50ep): mAP50-95 0,508 ← PEMENANG",
    "exp_02 (lr0.005, 50ep): 0,508",
    "exp_03 (lr0.003, 40ep): 0,473",
    "exp_04 (lr0.005, 30ep): 0,469",
    "Pemenang = konfigurasi baseline → tuning MENGONFIRMASI baseline optimal.",
], "presentation_figures/05_tune_01.png", "Perbandingan mAP50-95 antar eksperimen (pemenang disorot)",
   "Pemenang = baseline. Poin akademik: tuning tak selalu menaikkan skor — di sini mengonfirmasi konfigurasi awal sudah optimal. Hasil yang valid & jujur.")

figure_slide(22, "Tuning — Temuan: lr0 Diabaikan", [
    "exp_01 (lr0.01) & exp_02 (lr0.005) IDENTIK sampai 5 desimal.",
    "Sebab: optimizer='auto' menentukan lr sendiri → MENGABAIKAN lr0.",
    "Jadi variasi lr0 tak berefek; faktor nyata = jumlah epoch.",
    "Pelajaran: pahami perilaku framework (bukan asumsi).",
], "presentation_figures/05_tune_05.png", "Heatmap perbandingan metrik antar eksperimen",
   "Temuan analitis terkuat: dua lr berbeda hasil identik → lr0 diabaikan optimizer='auto'. Tunjukkan pemahaman mendalam framework — nilai plus untuk sidang.")

figure_slide(23, "Tuning — Pengaruh Jumlah Epoch", [
    "50 epoch > 40 > 30 (mAP50-95 0,508 > 0,473 > 0,469).",
    "Menambah epoch membantu konvergensi (sampai titik tertentu).",
    "lr0 tidak berefek (optimizer='auto') → epoch yang menentukan.",
], "presentation_figures/05_tune_04.png", "Metrik vs jumlah epoch",
   "Konfirmasi: yang berpengaruh adalah jumlah epoch, bukan lr0. Dasar saran: epoch lebih banyak / early-stopping.")

# ================== EVALUATION (DETAIL) ================== #
table_slide(24, "Evaluasi — Parameter & Alasan",
    ["Parameter", "Nilai", "Mengapa"],
    [["conf (threshold)", "0.25", "Titik OPERASI nyata (bukan 0.001 benchmark) → cerminkan penggunaan"],
     ["iou (NMS)", "0.5", "Ambang buang kotak tumpang-tindih saat prediksi"],
     ["imgsz", "640", "Sama dengan training (konsistensi)"],
     ["split", "val + test", "Ukur generalisasi (val vs test)"],
     ["sample prediksi", "12 gambar", "Verifikasi kualitatif hasil deteksi"]],
    "Evaluasi pada conf=0.25 (titik operasi) — ini sebabnya mAP eval lebih rendah dari mAP training (benchmark conf~0.001). iou=0.5 untuk NMS. Dievaluasi di val DAN test untuk cek generalisasi.",
    widths=[3.0, 2.0, 7.3])

table_slide(25, "Evaluasi — Val vs Test (Generalisasi)",
    ["Metrik", "Validation", "Test", "Selisih (val−test)"],
    [["Precision", f"{VAL['P']:.3f}", f"{TST['P']:.3f}", f"{VAL['P']-TST['P']:+.3f}"],
     ["Recall", f"{VAL['R']:.3f}", f"{TST['R']:.3f}", f"{VAL['R']-TST['R']:+.3f}"],
     ["mAP50", f"{VAL['m50']:.3f}", f"{TST['m50']:.3f}", f"{VAL['m50']-TST['m50']:+.3f}"],
     ["mAP50-95", f"{VAL['m5095']:.3f}", f"{TST['m5095']:.3f}", f"{VAL['m5095']-TST['m5095']:+.3f}"]],
    "Selisih mAP50-95 val vs test hanya ~0,047 (< 0,10) → generalisasi BAIK, tak ada overfitting parah (selaras EDA: split representatif). Catatan: test recall malah lebih tinggi, precision lebih rendah — model di test lebih 'berani'. mAP eval < mAP training karena conf=0.25 (titik operasi) vs benchmark.",
    widths=[3.0, 3.1, 3.1, 3.1],
    intro="Angka pada conf=0.25; mAP50-95 train (benchmark) = 0,508.")

figure_slide(26, "Evaluasi — Confusion Matrix", [
    "Pola kesalahan antar kelas pada validation.",
    "Diagonal = benar; off-diagonal = tertukar.",
    "No-Helmet paling menantang (minoritas + objek kecil).",
], "runs/evaluate/val/confusion_matrix.png", "Confusion matrix — evaluasi validation",
   "Per-kelas: kelas mana tertukar/terlewat. No-Helmet biasanya paling lemah (imbalance + objek kecil). Hubungkan ke EDA.")
figure_slide(27, "Evaluasi — Kurva Precision-Recall", [
    "Kurva PR per kelas: area di bawah = AP.",
    "Kelas dengan kurva lebih rendah = lebih sulit.",
    "Dasar memilih titik operasi (conf).",
], "runs/evaluate/val/BoxPR_curve.png", "Kurva Precision-Recall per kelas (validation)",
   "Kurva PR menunjukkan trade-off & AP per kelas. Kelas yang kurvanya 'jatuh' lebih cepat = lebih sulit (objek kecil/minoritas).")
figure_slide(28, "Evaluasi — Contoh Prediksi", [
    "Sample prediksi pada gambar validasi nyata.",
    "Kotak + label + confidence per objek.",
    "Verifikasi kualitatif kualitas deteksi.",
], "runs/evaluate/predictions/image0.jpg", "Contoh hasil prediksi (validation)",
   "Tunjukkan hasil deteksi nyata pada gambar. Diskusikan kasus benar & yang terlewat.")

# ================== DEPLOY ================== #
figure_slide(29, "Inference & Peringatan K3", [
    "Deteksi pada gambar uji: Helmet 8 · Person 7 · No-Helmet 1.",
    "Sistem memunculkan PERINGATAN K3 otomatis saat ada No-Helmet.",
    "Mendukung gambar / folder / video.",
], "runs/infer/predict/ppe_0855_png_jpg.rf.yQ34bQFHHa5jHuwcW3Es.jpg",
   "Contoh inference + peringatan K3 (terdeteksi 1 No-Helmet)",
   "Sisi aplikatif: model mendeteksi & memberi peringatan K3 saat ada pelanggaran. Angka nyata dari ringkasan inference. Tekankan nilai praktis untuk keselamatan.")

figure_slide(30, "Export & Deployment", [
    "Model diekspor ke ONNX (runs/export/best.onnx) — siap deploy.",
    "Paritas PyTorch vs ONNX terjaga → akurasi sama.",
    "Demo web: frontend (Vercel) + backend FastAPI (Hugging Face).",
    "Alternatif in-browser ONNX (tanpa server).",
], "presentation_figures/07_export_model_01.png", "Ukuran model antar-format",
   "Export ke ONNX untuk deployment (lepas dari PyTorch), akurasi terjaga. Sebutkan demo web yang dibangun (Vercel + HF FastAPI).")

# ================== PENUTUP ================== #
text_slide(31, "Diskusi & Keterbatasan", [
    "Recall moderat (~0,71): objek kecil & No-Helmet paling terdampak.",
    "Kualitas data: blur/low-light, label nyasar, potensi kebocoran antar-split.",
    "Tuning lr0 tak efektif (optimizer='auto') → faktor nyata = epoch.",
    "Metrik eval (conf 0.25) lebih rendah dari benchmark training — wajar (titik operasi).",
], "Jujur soal keterbatasan: recall belum tinggi, data menantang, tuning lr tak efektif karena framework. Penguji menghargai kesadaran kritis & pemahaman batas.")

text_slide(32, "Kesimpulan", [
    "Pipeline deteksi helm 3 kelas yang LENGKAP & reproducible berhasil dibangun (8 tahap).",
    f"Model baseline yolo26s — training mAP50-95 {TR['m5095']:.3f} (best epoch 40);",
    (f"evaluasi (conf 0.25): val {VAL['m5095']:.3f} · test {TST['m5095']:.3f} → generalisasi baik (selisih ~0,047).", 1),
    "Karakter: precision tinggi (~0,85), recall moderat (~0,71) → konservatif.",
    "Tantangan utama (objek kecil & No-Helmet) konsisten dari EDA hingga evaluasi.",
    "Tuning mengonfirmasi konfigurasi baseline optimal; temuan lr0 diabaikan optimizer='auto'.",
    "Setiap keputusan parameter dijustifikasi data; model diekspor (ONNX) + demo web siap pakai.",
], "Kesimpulan menyeluruh: pipeline lengkap berhasil, model bekerja & menggeneralisasi, tiap keputusan berbasis data. Tekankan benang merah masalah (objek kecil/No-Helmet) dari EDA→evaluasi, dan kontribusi reproducibility + deployment.")

text_slide(33, "Saran & Pekerjaan Mendatang", [
    "Naikkan recall: imgsz lebih tinggi (768) untuk objek kecil; tambah data No-Helmet.",
    "Augmentasi terarah untuk kelas minoritas / objek kecil.",
    "Set optimizer eksplisit (SGD/AdamW) agar tuning learning rate efektif.",
    "Bersihkan/periksa kebocoran data antar-split untuk evaluasi lebih jujur.",
    "Deploy nyata (HF Space + Vercel) & uji pada rekaman CCTV lapangan.",
], "Arah perbaikan konkret berbasis temuan: resolusi lebih tinggi, lebih banyak data No-Helmet, optimizer eksplisit, atasi kebocoran data, lalu deployment & uji lapangan.")

cover("Terima Kasih", "Diskusi & Tanya Jawab",
      ["Siap menjawab: recall rendah · imbalance · optimizer='auto' · rencana deployment"],
      "Penutup & Q&A. Siapkan jawaban untuk pertanyaan umum: kenapa recall rendah, cara atasi imbalance, kenapa optimizer='auto' penting, rencana deployment.")

# ----------------------------- Simpan -------------------------------- #
pptx_path = OUT / "Presentasi_Helmet_YOLO26.pptx"
prs.save(pptx_path)
print("PPTX:", pptx_path, "|", len(prs.slides._sldIdLst), "slide")

md = ["# Script Presentasi — Deteksi Helm Konstruksi (YOLO26s)\n",
      "Narasi per slide untuk sidang/akademik (data nyata terkini).\n"]
for num, title, items, note in SLIDES:
    md.append(f"\n## Slide {num}: {title}\n\n**Poin di layar:**\n")
    for b in items:
        md.append(("  - " + b[0]) if isinstance(b, tuple) else ("- " + str(b)))
    md.append("\n**Narasi:**\n\n" + note + "\n")
(OUT / "Script_Presentasi.md").write_text("\n".join(md), encoding="utf-8")
print("Script:", OUT / "Script_Presentasi.md")
